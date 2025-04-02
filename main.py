from tercen.client import context as ctx
import numpy as np
import struct
from hashlib import md5
import polars as pl
from base64 import b64encode
from pptx import Presentation
from pptx.util import Inches
import os

tercenCtx = ctx.TercenContext(
    workflowId="241160841cead65995c0c45d0a011ce8",
    stepId="0e756662-c900-41f8-bde2-f5917d8d7779",
    username="admin", # if using the local Tercen instance
    password="admin", # if using the local Tercen instance
    serviceUri = "http://127.0.0.1:5402/" # if using the local Tercen instance 
)

data = tercenCtx.select(['.y', '.ci', '.ri'], df_lib="pandas").values
data = np.array(data, dtype=np.uint32)

R = (data[:, 0] >> 16) & 0xFF
G = (data[:, 0] >> 8) & 0xFF
B = data[:, 0] & 0xFF

X = data[:, 1]
Y = data[:, 2]

pixels = np.column_stack((R, G, B, X, Y))

def save_as_bmp(filename, pixel_data):
    if not isinstance(pixel_data, np.ndarray) or pixel_data.shape[1] != 5:
        raise ValueError("Pixel data must be a NumPy array with shape (n, 5) for [R,G,B,X,Y]")
    
    width = int(pixel_data[:, 3].max() + 1)
    height = int(pixel_data[:, 4].max() + 1)
    padding = (4 - (width * 3) % 4) % 4
    row_size = (width * 3) + padding
    
    file_size = 14 + 40 + (row_size * height)
    file_header = struct.pack('<2sIHHI', b'BM', file_size, 0, 0, 54)
    dib_header = struct.pack('<IIIHHIIIIII', 40, width, height, 1, 24, 0, row_size * height, 2835, 2835, 0, 0)
    
    image_data = bytearray([0] * (row_size * height))
    for pixel in pixel_data:
        r, g, b, x, y = pixel
        x, y = int(x), int(y)
        if 0 <= x < width and 0 <= y < height:
            pos = (height - 1 - y) * row_size + x * 3
            image_data[pos:pos+3] = [int(b), int(g), int(r)]
    
    with open(filename, 'wb') as f:
        f.write(file_header)
        f.write(dib_header)
        f.write(image_data)

save_as_bmp('output.bmp', pixels)




img_path = "output.bmp"


prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
height = Inches(5.5)

pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs.save('test.pptx')



with open("test.pptx", "rb") as file:
    fileBytes = file.read()
    checksum = md5(fileBytes).hexdigest()

    imgDf = pl.DataFrame({\
        ".ci":[0],\
        "filename":["compressed image"],\
        "mimetype":["application/vnd.ms-powerpoint"],\
        "checksum":[checksum],\
        ".content":[str(b64encode(fileBytes).decode("utf-8"))]\
        })
    
    imgDf = imgDf.with_columns(pl.col('.ci').cast(pl.Int32))
    df = tercenCtx.add_namespace(imgDf) 

    os.remove("output.bmp")
    os.remove("test.pptx")

    tercenCtx.save(df)


